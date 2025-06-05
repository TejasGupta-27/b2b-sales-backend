import os
import json
from openai import AzureOpenAI
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml


# Initialize Azure OpenAI client
client = AzureOpenAI(
    api_key=os.getenv("AZURE_OPENAI_API_KEY") or "***REMOVED***",
    api_version="2025-01-01-preview",
    azure_endpoint="https://internship2025-teama.openai.azure.com"
)

DEPLOYMENT_NAME = "gpt-4.1-mini"

def hide_placeholders(slide):
    """Safely hide placeholders without corrupting the slide structure"""
    shapes_to_remove = []
    
    for shape in slide.shapes:
        # Check if it's a placeholder
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
            shapes_to_remove.append(shape)
        # Also remove text boxes that might be default placeholders
        elif hasattr(shape, 'text_frame') and shape.text_frame is not None:
            # Check if it contains placeholder text
            if shape.text_frame.text.strip() in ['Click to add title', 'Click to add subtitle', 'Click to add text']:
                shapes_to_remove.append(shape)
    
    # Remove identified placeholders
    for shape in shapes_to_remove:
        try:
            sp = shape._element
            sp.getparent().remove(sp)
        except:
            # If direct removal fails, try making it invisible
            try:
                shape.width = 0
                shape.height = 0
            except:
                pass

def clear_slide_content_safely(slide):
    """Safely clear slide content without corrupting the slide structure"""
    # Only remove content shapes, not structural elements
    shapes_to_remove = []
    
    for shape in slide.shapes:
        # Only remove shapes that are not essential slide structure
        if (hasattr(shape, 'shape_type') and 
            shape.shape_type in [1, 17, 14]):  # Text box, auto shape, picture
            shapes_to_remove.append(shape)
        elif hasattr(shape, 'text_frame') and shape.text_frame is not None:
            # Clear text content but keep the shape structure
            try:
                shape.text_frame.clear()
            except:
                shapes_to_remove.append(shape)
    
    # Remove non-essential shapes
    for shape in shapes_to_remove:
        try:
            sp = shape._element
            sp.getparent().remove(sp)
        except:
            # If removal fails, just hide it
            try:
                shape.width = 0
                shape.height = 0
            except:
                pass


def extract_ppt_structure(quotation: str) -> dict:
    """Use Azure OpenAI to generate a detailed and persuasive sales pitch deck structure from the quotation."""
    prompt = f"""
You are a business assistant. Based on the product quotation below, generate a structured and persuasive PowerPoint sales pitch deck in **valid JSON** format.

### QUOTATION
\"\"\"
{quotation}
\"\"\"

### TASKS
1. Analyze the quotation to identify:
   - Customer name
   - Product name
   - Specifications (CPU, RAM, Storage, etc.)
   - Price, Delivery Timeline, Warranty, Support options

2. Generate a slide deck in this order:
   1. Customer Need
   2. Our Solution
   3. Product Overview (specs)
   4. Pricing Breakdown
   5. Warranty & Support
   6. Product Comparison (see below)
   7. Delivery Timeline
   8. Call to Action

Each slide must contain a **title** and 5–6 persuasive bullet points.

3. Add a **comparison table** duplicating the same product 3 times with slightly varied names:
   - Copy the product specs from the quotation
   - Change the name to `"Product Name Variant 1"`, `"Variant 2"`, etc.
   - Add them to a table with the following structure:

### JSON OUTPUT FORMAT
Return your response as valid JSON:
{{
  "slides": [
    {{
      "title": "Slide Title",
      "content": ["Bullet 1", "Bullet 2", "..."]
    }},
    ...
  ],
  "tables": [
    {{
      "title": "Product Comparison",
      "columns": ["Product Name", "Price", "CPU", "RAM", "Storage", "Warranty", "Support"],
      "rows": [
        ["...Variant 1...", "...", "...", "...", "...", "...", "..."],
        ["...Variant 2...", "...", "...", "...", "...", "...", "..."],
        ["...Variant 3...", "...", "...", "...", "...", "...", "..."]
      ]
    }}
  ]
}}

✅ Use ONLY the product and specifications mentioned in the quotation — do NOT make up new ones.  
✅ Use slightly varied product names for realism.  
✅ Return valid JSON ONLY — no commentary, no markdown.
"""

    response = client.chat.completions.create(
        model=DEPLOYMENT_NAME,
        messages=[{"role": "user", "content": prompt}],
        temperature=0.5
    )

    raw_output = response.choices[0].message.content.strip()

    try:
        parsed = json.loads(raw_output)
        if isinstance(parsed, list):
            return {"slides": parsed}
        elif isinstance(parsed, dict) and "slides" in parsed:
            return parsed
        else:
            raise ValueError("Parsed JSON is missing expected 'slides' structure.")
    except json.JSONDecodeError:
        print("❌ GPT response was not valid JSON. Response was:\n", raw_output)
        raise


def add_comparison_table(slide, table_data):
    rows = len(table_data["rows"]) + 1  # +1 for header
    cols = len(table_data["columns"])
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Header row
    for col, header in enumerate(table_data["columns"]):
        cell = table_shape.cell(0, col)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)

    # Data rows
    for i, row in enumerate(table_data["rows"], start=1):
        for j, value in enumerate(row):
            cell = table_shape.cell(i, j)
            cell.text = str(value)  # Ensure it's a string
            cell.text_frame.paragraphs[0].font.size = Pt(12)


from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor

prs = Presentation()

LOGO_PATH = "assets/logo.png"

def add_logo(slide):
    try:
        slide.shapes.add_picture(LOGO_PATH, Inches(8.5), Inches(0.2), width=Inches(1.2))
    except Exception as e:
        print("⚠️ Logo not added:", e)

def add_price_chart(slide, product_name, price):
    chart_data = CategoryChartData()
    chart_data.categories = [product_name, "Competitor A", "Competitor B"]
    chart_data.add_series('Price ($)', [price, price * 1.2, price * 1.1])

    x, y, cx, cy = Inches(0.5), Inches(2), Inches(8), Inches(4)
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data)

def add_contextual_image(slide, title_text):
    image_path = None
    title = title_text.lower()
    if "need" in title:
        image_path = "assets/problem.png"
    elif "solution" in title:
        image_path = "assets/solution.png"
    elif "support" in title:
        image_path = "assets/support.png"
    elif "overview" in title:
        image_path = "assets/specs.png"
    elif "comparison" in title:
        image_path = "assets/compare.png"
    elif "call to action" in title:
        image_path = "assets/action.png"

    if image_path:
        try:
            slide.shapes.add_picture(image_path, Inches(0.2), Inches(1.5), width=Inches(1.0))
        except Exception as e:
            print("⚠️ Image not added:", e)

def set_background_color(slide, rgb_color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb_color

def add_timeline_shape(slide):
    for i in range(3):
        left = Inches(1 + i * 2.5)
        top = Inches(2)
        box = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(1), Inches(1))
        box.text = f"Step {i+1}"

def generate_ppt(presentation_structure, product_name, product_price):
    for slide_data in presentation_structure:
        slide_layout = prs.slide_layouts[5]  # Title Only
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        title.text = slide_data["title"]

        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5.5)
        text_box = slide.shapes.add_textbox(left, top, width, height)
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data["content"]
        p.font.size = Pt(20)

        # Enhancements
        add_logo(slide)
        set_background_color(slide, RGBColor(245, 245, 245))
        add_contextual_image(slide, slide_data["title"])

        if "price" in slide_data["title"].lower():
            add_price_chart(slide, product_name, product_price)
        elif "timeline" in slide_data["title"].lower():
            add_timeline_shape(slide)

    prs.save("output/enhanced_presentation.pptx")

    # Save presentation
    try:
        prs.save(output_path)
        print(f"✅ Presentation saved successfully to: {output_path}")
    except Exception as e:
        print(f"❌ Error saving presentation: {e}")
        raise


if __name__ == "__main__":
    quotation_input = """
    Here is the formal quote for one unit of the Workstation Pro Professional:

    Customer: [Your Company Name]
    Product: Workstation Pro Professional

    CPU: Intel Xeon W-2295
    RAM: 32GB DDR4
    Storage: 1TB NVMe SSD
    Quantity: 1 unit
    Unit Price: $3,499.99
    Delivery Timeline: 2-3 weeks
    Warranty: 1-year included
    Support: Optional setup and integration
    """

    try:
        print("🧠 Generating structured pitch deck from quotation...")
        structured_data = extract_ppt_structure(quotation_input)

        print("🖼️ Creating PowerPoint presentation...")
        generate_ppt(structured_data)

    except Exception as e:
        print("❌ Error:", e)