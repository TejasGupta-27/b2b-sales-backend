import os
import json
from openai import AzureOpenAI
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
import logging
from typing import List, Dict, Any
from services.localisation import get_category_translation

logger = logging.getLogger(__name__)

# --- DESIGN CONSTANTS ---
COVER_TITLE_FONT_SIZE = Pt(54)
COVER_TITLE_BOX = (Inches(1), Inches(1.2), Inches(8), Inches(1.5))
COVER_TITLE_MARGIN = (Inches(0.3), Inches(0.3), Inches(0.2), Inches(0.2))
COVER_TOP_LINE = (Inches(1.5), Inches(1.0), Inches(7), Inches(0.12))
COVER_BOTTOM_LINE = (Inches(1.5), Inches(2.5), Inches(7), Inches(0.12))
COVER_SUBTITLE_BOX = (Inches(1), Inches(3.2), Inches(8), Inches(1))
COVER_SUBTITLE_MARGIN = (Inches(0.3), Inches(0.3), Inches(0.1), Inches(0.1))
COVER_SUBTITLE_FONT_SIZE = Pt(28)
COVER_LOGO_BOTTOM_LEFT = (Inches(0.2), Inches(6.5), Inches(1.0))
COVER_LOGO_TOP_LEFT = (Inches(0.2), Inches(0.2), Inches(1.0))
COVER_IMAGE_COVER = (Inches(2.5), Inches(4.5), Inches(4.0))
COVER_IMAGE_NORMAL = (Inches(7.2), Inches(5.7), Inches(2.0))

CONTENT_TITLE_BOX = (Inches(0.5), Inches(0.3), Inches(9), Inches(1.2))
CONTENT_TITLE_MARGIN = (Inches(0.2), Inches(0.2), Inches(0.1), Inches(0.1))
CONTENT_TITLE_FONT_SIZE = Pt(38)
CONTENT_UNDERLINE = (Inches(2), Inches(1.3), Inches(6), Inches(0.07))
CONTENT_BOX = (Inches(0.8), Inches(1.8), Inches(8.4), Inches(5.5))
CONTENT_MARGIN = (Inches(0.3), Inches(0.3), Inches(0.2), Inches(0.2))
CONTENT_FONT_SIZE = Pt(20)
CONTENT_SLIDE_NUMBER_BOX = (Inches(8.5), Inches(7.2), Inches(1), Inches(0.5))
CONTENT_SLIDE_NUMBER_FONT_SIZE = Pt(12)

TABLE_BOX = (Inches(0.5), Inches(1.8), Inches(9), Inches(4.5))
TABLE_HEADER_FONT_SIZE = Pt(16)
TABLE_HEADER_MARGIN = (Inches(0.1), Inches(0.1), Inches(0.1), Inches(0.1))
TABLE_ROW_FONT_SIZE = Pt(14)
TABLE_ROW_MARGIN = (Inches(0.1), Inches(0.1), Inches(0.1), Inches(0.1))

class PitchDeckService:
    def __init__(self):
        # Initialize Azure OpenAI client - disabled for local testing
        self.client_configured = False
        self.deployment_name = "gpt-4"  # Default deployment name for when enabled
        
        # For production, uncomment this block:
        
        from config import settings
        try:
            self.client_configured = True
            self.client = AzureOpenAI(
                api_key=settings.azure_openai_api_key,
                api_version=settings.azure_openai_api_version,
                azure_endpoint=settings.azure_openai_endpoint
            )
            self.deployment_name = settings.azure_openai_deployment_name
        except Exception as e:
            logger.error(f"❌ Failed to initialize Azure OpenAI client: {e}")
            self.client_configured = False
        
        
        # Font configuration for Japanese support
        self.japanese_fonts = [
            "Noto Sans CJK JP",  # Primary choice - commonly available
            "Yu Gothic UI",       # Windows default
            "MS Gothic",          # Common fallback
            "Hiragino Kaku Gothic ProN",
            "Hiragino Sans",
            "Arial Unicode MS",   # Last resort fallback
        ]
        
        # Detect if system supports Japanese fonts
        self.title_font = self._get_available_font(bold=True)
        self.body_font = self._get_available_font(bold=False)
        
        logger.info(f"Using fonts - Title: {self.title_font}, Body: {self.body_font}")
    
    def _get_available_font(self, bold=False):
        """Get the best available font for Japanese text"""
        # For PowerPoint, prioritize Noto Sans CJK JP if available
        preferred_fonts = []
        
        if bold:
            # For titles, try bold variants first
            preferred_fonts = [
                "Noto Sans CJK JP Bold",
                "Yu Gothic UI Semibold",
                "MS Gothic Bold",
            ] + self.japanese_fonts
        else:
            # For body text, use regular variants
            preferred_fonts = self.japanese_fonts
        
        # Return Noto Sans as our primary choice (we know it's available)
        return "Noto Sans CJK JP"
    
    def _detect_japanese_text(self, text):
        """Detect if text contains Japanese characters"""
        if not text:
            return False
        
        # Check for Japanese character ranges
        for char in text:
            code = ord(char)
            # Hiragana: 0x3040-0x309F
            # Katakana: 0x30A0-0x30FF  
            # Kanji: 0x4E00-0x9FAF
            # Japanese punctuation: 0x3000-0x303F
            if (0x3000 <= code <= 0x303F or  # Japanese punctuation
                0x3040 <= code <= 0x309F or  # Hiragana
                0x30A0 <= code <= 0x30FF or  # Katakana
                0x4E00 <= code <= 0x9FAF):   # Kanji
                return True
        return False
    
    def _apply_font_to_paragraph(self, paragraph, text, is_title=False):
        """Apply appropriate font to paragraph based on text content"""
        paragraph.text = text
        
        # Detect if Japanese text is present
        has_japanese = self._detect_japanese_text(text)
        
        if has_japanese:
            # Use Japanese-compatible font
            font_name = self.title_font if is_title else self.body_font
        else:
            # Use standard fonts for English text
            font_name = "Segoe UI Semibold" if is_title else "Segoe UI"
        
        # Apply font settings
        paragraph.font.name = font_name
        
        # Set font for Asian text specifically (this helps with mixed content)
        if has_japanese:
            try:
                # This helps PowerPoint use the correct font for Asian characters
                paragraph.font._element.set(qn('a:eastAsianTheme'), 'minor')
            except:
                pass  # Ignore if this advanced setting fails
        
        return paragraph

    def hide_placeholders(self, slide):
        shapes_to_remove = []
        for shape in slide.shapes:
            if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
                shapes_to_remove.append(shape)
            elif hasattr(shape, 'text_frame') and shape.text_frame is not None:
                if shape.text_frame.text.strip() in ['Click to add title', 'Click to add subtitle', 'Click to add text']:
                    shapes_to_remove.append(shape)
        for shape in shapes_to_remove:
            try:
                sp = shape._element
                sp.getparent().remove(sp)
            except:
                try:
                    shape.width = 0
                    shape.height = 0
                except:
                    pass

    def clear_slide_content_safely(self, slide):
        shapes_to_remove = []
        for shape in slide.shapes:
            if hasattr(shape, 'shape_type') and shape.shape_type in [1, 17, 14]:
                shapes_to_remove.append(shape)
            elif hasattr(shape, 'text_frame') and shape.text_frame is not None:
                try:
                    shape.text_frame.clear()
                except:
                    shapes_to_remove.append(shape)
        for shape in shapes_to_remove:
            try:
                sp = shape._element
                sp.getparent().remove(sp)
            except:
                try:
                    shape.width = 0
                    shape.height = 0
                except:
                    pass

    # ✅ UPDATED: improved language detection and hybrid fallback
    async def extract_ppt_structure(self, quotation: str, language: str = "en", include_comparison_table: bool = False) -> dict:
        """Generate a structured and persuasive sales pitch deck with hybrid language detection."""
        
        # Initialize language resolution with explicit setting
        resolved_language = language if language else "en"
        detection_method = "explicit" if language else "default"
        confidence = 1.0
        
        # If no explicit language provided, try to detect Japanese
        if not language:
            has_japanese = any(self._detect_japanese_text(line) for line in quotation.split('\n') if line.strip())
            if has_japanese:
                resolved_language = "ja"
                detection_method = "character_detection"
                confidence = 0.95
                logger.info(f"🔍 Detected Japanese text in the quotation")
        
        # Package language resolution info
        language_resolution = {
            "language": resolved_language,
            "method": detection_method,
            "confidence": confidence
        }
        
        logger.info(f"🌐 Language Resolution: {resolved_language} (method: {detection_method}, confidence: {confidence:.2f})")
        
        # Use Azure OpenAI if configured, otherwise use fallback
        if not self.client_configured:
            return self._get_fallback_structure(resolved_language)
        
        # Hybrid language detection approach
        from services.language_service import LanguageService
        language_service = LanguageService()
        
        # Ensure we respect the input language first, then try detection
        if language and language != "en":
            resolved_language = language
            detection_method = "explicit"
            confidence = 1.0
        else:
            # Fallback to hybrid approach
            language_resolution = language_service.resolve_language(
                explicit_language=language,
                text_content=quotation,
                context=None
            )
            resolved_language = language_resolution['language']
            detection_method = language_resolution['method']
            confidence = language_resolution['confidence']
        
        logger.info(f"🌐 Pitch Deck Language Resolution:")
        logger.info(f"   Input Language: {language}")
        logger.info(f"   Resolved Language: {resolved_language}")
        logger.info(f"   Detection Method: {detection_method}")
        logger.info(f"   Confidence: {confidence:.2f}")

        # Add language condition based on resolved language
        language_note = ""
        if resolved_language == "ja":
            language_note = "Ensure output is in Japanese (日本語). Use natural and business-appropriate Japanese expressions."
        elif resolved_language == "es":
            language_note = "Ensure output is in Spanish. Use natural and business-appropriate Spanish expressions."
        elif resolved_language == "fr":
            language_note = "Ensure output is in French. Use natural and business-appropriate French expressions."
        elif resolved_language == "de":
            language_note = "Ensure output is in German. Use natural and business-appropriate German expressions."
        elif resolved_language != "en":
            language_note = f"Ensure output is in {language_service.get_language_name(resolved_language)}."

        prompt = f"""
You are a business assistant. Based on the product quotation below, generate a structured and persuasive PowerPoint sales pitch deck in **valid JSON** format.

### QUOTATION
\"""
{quotation}
\"""

### TASKS
1. Analyze the quotation to identify:
   - Customer name
   - Product name
   - Specifications (CPU, RAM, Storage, etc.)
   - Price, Delivery Timeline, Warranty, Support options

2. Generate a slide deck in this order:
   1. Customer Need
   2. Our Solution
   3. Product Overview (specs)
   4. Pricing Breakdown
   5. Warranty & Support
   6. Delivery Timeline

Each slide must contain a **title** and 5–6 persuasive bullet points.

### JSON OUTPUT FORMAT
Return your response as valid JSON:
{{
  "slides": [
    {{
      "title": "Slide Title",
      "content": ["Bullet 1", "Bullet 2", "..."]
    }},
    ...
  ]
}}

✅ Use ONLY the product and specifications mentioned in the quotation — do NOT make up new ones.
✅ Return valid JSON ONLY — no commentary, no markdown.

{language_note}
"""
        
        # If comparison table is requested, add it to the prompt
        if include_comparison_table:
            comparison_section = """
3. Add a **comparison table** for competitive analysis with the following structure:

Add this to your JSON response:
"tables": [
  {
    "title": "Product Comparison",
    "columns": ["Product Name", "Price", "CPU", "RAM", "Storage", "Warranty", "Support"],
    "rows": [
      ["Product 1", "Price 1", "CPU 1", "RAM 1", "Storage 1", "Warranty 1", "Support 1"],
      ["Product 2", "Price 2", "CPU 2", "RAM 2", "Storage 2", "Warranty 2", "Support 2"],
      ["Product 3", "Price 3", "CPU 3", "RAM 3", "Storage 3", "Warranty 3", "Support 3"]
    ]
  }
]

Note: The comparison table will be populated with real competitor data separately.
"""
            # Insert the comparison section before the JSON format
            prompt = prompt.replace("### JSON OUTPUT FORMAT", comparison_section + "\n### JSON OUTPUT FORMAT")

        try:
            response = self.client.chat.completions.create(
                model=self.deployment_name,
                messages=[{"role": "user", "content": prompt}],
                temperature=0.5
            )

            raw_output = response.choices[0].message.content.strip()

            try:
                parsed = json.loads(raw_output)
                if isinstance(parsed, list):
                    result = {"slides": parsed}
                elif isinstance(parsed, dict) and "slides" in parsed:
                    result = parsed
                else:
                    raise ValueError("Parsed JSON is missing expected 'slides' structure.")
                
                # Add language metadata to result
                result['language_resolution'] = language_resolution
                result['resolved_language'] = resolved_language
                
                return result
                
            except json.JSONDecodeError:
                logger.error("GPT response was not valid JSON. Response was:\n%s", raw_output)
                return self._get_fallback_structure(resolved_language)
                
        except Exception as e:
            logger.error(f"Error in extract_ppt_structure: {e}")
            return self._get_fallback_structure(resolved_language)
    
    def _get_fallback_structure(self, language: str = "en") -> dict:
        """Provide a fallback structure when OpenAI is not available with language support"""
        
        # Create language resolution metadata
        language_resolution = {
            "language": language,
            "method": "fallback",
            "confidence": 1.0
        }
        
        # Get localized fallback content
        if language == "ja":
            return {
                "slides": [
                    {
                        "title": "お客様のニーズ分析",
                        "content": [
                            "ビジネス要件の理解",
                            "技術的な課題と機会の特定",
                            "現在のインフラの制約評価",
                            "スケーラビリティと成長ニーズの決定",
                            "予算とタイムラインの制約評価"
                        ]
                    },
                    {
                        "title": "包括的なソリューション",
                        "content": [
                            "お客様のビジネスに合わせたテクノロジーソリューション",
                            "業界をリードする製品とサービス",
                            "既存システムとのシームレスな統合",
                            "包括的なサポートとメンテナンス",
                            "将来に対応したスケーラブルなアーキテクチャ"
                        ]
                    },
                    {
                        "title": "製品概要と仕様",
                        "content": [
                            "高性能ハードウェアコンポーネント",
                            "エンタープライズグレードソフトウェアソリューション",
                            "高度なセキュリティ機能とプロトコル",
                            "信頼性の高いパフォーマンスと稼働時間保証",
                            "エネルギー効率的で持続可能な設計"
                        ]
                    },
                    {
                        "title": "投資と価格の内訳",
                        "content": [
                            "プレミアムソリューションの競争力のある価格",
                            "隠れた料金のない透明なコスト構造",
                            "柔軟な支払いオプションと融資",
                            "大量購入の場合のボリュームディスカウント",
                            "費用対効果の高い総所有コストモデル"
                        ]
                    },
                    {
                        "title": "保証とサポートサービス",
                        "content": [
                            "包括的な保証カバレッジ",
                            "24時間365日の技術サポート体制",
                            "オンサイトサービスとメンテナンス",
                            "リモート監視と診断",
                            "定期的なアップデートとシステム最適化"
                        ]
                    },
                    {
                        "title": "配送スケジュールと実装",
                        "content": [
                            "効率的なプロジェクト計画と実行",
                            "日常業務への最小限の影響",
                            "段階的な実装アプローチ",
                            "スタッフトレーニングと知識移転",
                            "実装後のサポートと最適化"
                        ]
                    },
                    {
                        "title": "次のステップとお取引",
                        "content": [
                            "詳細な技術コンサルテーションのスケジューリング",
                            "ソリューション仕様の確認と最終化",
                            "プロジェクトタイムラインとマイルストーンの確認",
                            "契約書に署名し、実装を開始",
                            "技術エクセレンスへの旅を始めましょう"
                        ]
                    }
                ],
                "language_resolution": {
                    "language": language,
                    "method": "fallback",
                    "confidence": 1.0
                },
                "resolved_language": language
            }
        else:
            # English fallback (existing structure)
            return {
                "slides": [
                    {
                        "title": "Customer Needs Analysis",
                        "content": [
                            "Understanding your business requirements",
                            "Identifying technology gaps and opportunities",
                            "Assessing current infrastructure limitations",
                            "Determining scalability and growth needs",
                            "Evaluating budget and timeline constraints"
                        ]
                    },
                    {
                        "title": "Our Comprehensive Solution",
                        "content": [
                            "Tailored technology solutions for your business",
                            "Industry-leading products and services",
                            "Seamless integration with existing systems",
                            "Comprehensive support and maintenance",
                            "Future-ready scalable architecture"
                        ]
                    },
                    {
                        "title": "Product Overview & Specifications",
                        "content": [
                            "High-performance hardware components",
                            "Enterprise-grade software solutions",
                            "Advanced security features and protocols",
                            "Reliable performance and uptime guarantees",
                            "Energy-efficient and sustainable design"
                        ]
                    },
                    {
                        "title": "Investment & Pricing Breakdown",
                        "content": [
                            "Competitive pricing for premium solutions",
                            "Transparent cost structure with no hidden fees",
                            "Flexible payment options and financing",
                            "Volume discounts for bulk purchases",
                            "Cost-effective total ownership model"
                        ]
                    },
                    {
                        "title": "Warranty & Support Services",
                        "content": [
                            "Comprehensive warranty coverage",
                            "24/7 technical support availability",
                            "On-site service and maintenance",
                            "Remote monitoring and diagnostics",
                            "Regular updates and system optimization"
                        ]
                    },
                    {
                        "title": "Delivery Timeline & Implementation",
                        "content": [
                            "Efficient project planning and execution",
                            "Minimal disruption to daily operations",
                            "Phased implementation approach",
                            "Staff training and knowledge transfer",
                            "Post-implementation support and optimization"
                        ]
                    }
                ],
                "language_resolution": {
                    "language": language,
                    "method": "fallback",
                    "confidence": 1.0
                },
                "resolved_language": language
            }

    def add_comparison_table(self, slide, table_data):
        """Add comparison table with improved styling and Japanese font support"""
        rows = len(table_data["rows"]) + 1  # +1 for header
        cols = len(table_data["columns"])
        left, top, width, height = TABLE_BOX
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table
        # Set table style
        table_shape.table_direction = 0  # Left to right
        # Add black borders to all cells (namespace safe)
        NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
        nsmap = {'a': NS_A}
        for row in range(rows):
            for col in range(cols):
                cell = table_shape.cell(row, col)
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                for border in ["lnL", "lnR", "lnT", "lnB"]:
                    tag = f'{{{NS_A}}}{border}'
                    ln = tcPr.find(tag)
                    if ln is None:
                        ln = tcPr.makeelement(tag)
                        tcPr.append(ln)
                    ln.set("w", "12700")  # 1pt
                    ln.set("cap", "flat")
                    ln.set("cmpd", "sng")
                    ln.set("algn", "ctr")
                    # Set solidFill for black
                    solidFill = ln.find(f'{{{NS_A}}}solidFill')
                    if solidFill is None:
                        solidFill = ln.makeelement(f'{{{NS_A}}}solidFill')
                        ln.append(solidFill)
                    srgbClr = solidFill.find(f'{{{NS_A}}}srgbClr')
                    if srgbClr is None:
                        srgbClr = ln.makeelement(f'{{{NS_A}}}srgbClr')
                        solidFill.append(srgbClr)
                    srgbClr.set("val", "000000")

        # Header row with improved styling
        for col, header in enumerate(table_data["columns"]):
            cell = table_shape.cell(0, col)
            cell.text_frame.clear()
            p = cell.text_frame.paragraphs[0]
            self._apply_font_to_paragraph(p, header, is_title=True)
            p.font.bold = True
            p.font.size = TABLE_HEADER_FONT_SIZE
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(231, 76, 60)
            cell.margin_left, cell.margin_right, cell.margin_top, cell.margin_bottom = TABLE_HEADER_MARGIN
        for i, row in enumerate(table_data["rows"], start=1):
            for j, value in enumerate(row):
                cell = table_shape.cell(i, j)
                cell.text_frame.clear()
                p = cell.text_frame.paragraphs[0]
                self._apply_font_to_paragraph(p, str(value), is_title=False)
                p.font.size = TABLE_ROW_FONT_SIZE
                p.font.color.rgb = RGBColor(44, 62, 80)
                p.alignment = PP_ALIGN.CENTER
                if i % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(248, 249, 250)
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                cell.margin_left, cell.margin_right, cell.margin_top, cell.margin_bottom = TABLE_ROW_MARGIN

    def create_comparison_table_from_products(self, similar_products: List[Dict[str, Any]], title: str = "Product Comparison", language: str = "en") -> Dict[str, Any]:
        """Create a comparison table with proper localization"""
        if not similar_products:
            logger.warning("No similar products provided for comparison table")
            return None
        
        # Localized headers
        headers = {
            "ja": ["製品名", "主な機能", "価格", "販売元"],
            "en": ["Product Name", "Key Features", "Price", "Vendor"]
        }.get(language, ["Product Name", "Key Features", "Price", "Vendor"])
        
        # Currency symbols
        currency_symbols = {
            "ja": "¥",
            "en": "$"
        }
        currency_symbol = currency_symbols.get(language, "$")
        
        # Define table structure
        comparison_table = {
            "title": title,
            "columns": headers,
            "rows": []
        }
        
        # Process up to 3 products
        for i, product in enumerate(similar_products[:3]):
            # Format price based on language
            price = product.get('price', 0)
            if language == "ja":
                price_str = f"{currency_symbol}{int(price * 150):,}"  # Rough USD to JPY conversion
            else:
                price_str = f"{currency_symbol}{price:,.2f}"
            
            # Get localized name if available
            name = product.get(f'name_{language}', product.get('name', 'Unknown'))
            
            # Get localized features
            features = product.get(f'features_{language}', product.get('features', []))
            if isinstance(features, list):
                feature_text = " • ".join(features[:3])  # Top 3 features
            else:
                feature_text = str(features)
            
            # Get localized vendor
            vendor = product.get(f'vendor_{language}', product.get('vendor', 'Unknown'))
            
            row = [name, feature_text, price_str, vendor]
            comparison_table["rows"].append(row)
        
        # Ensure we have at least 3 rows for better presentation
        while len(comparison_table["rows"]) < 3:
            empty_row = ["--"] * len(headers)
            comparison_table["rows"].append(empty_row)
        
        logger.info(f"Created comparison table with {len(comparison_table['rows'])} products in {language}")
        return comparison_table

    def get_product_category(self, product_category: str, language: str = "en") -> dict:
        """
        Returns a dict with both the English and localized product category names, falling back to 'general products' if not found.
        { 'en': <english_name>, 'localized': <localized_name> }
        """
        import warnings
        from services.localisation import get_category_translation
        fallback_en = "general products"
        fallback_localized = get_category_translation("general products", language)
        if not product_category:
            warnings.warn(f"No category provided, falling back to '{fallback_localized}'")
            return {"en": fallback_en, "localized": fallback_localized}
        # Check if category file exists
        base_dir = os.path.dirname(os.path.abspath(__file__))
        category_file = os.path.abspath(os.path.join(base_dir, "..", "Data", "json", f"{product_category}.json"))
        if not os.path.exists(category_file):
            warnings.warn(f"Category file '{category_file}' not found, falling back to '{fallback_localized}'")
            return {"en": fallback_en, "localized": fallback_localized}
        # Return both English and localized category name
        return {"en": product_category, "localized": get_category_translation(product_category, language)}

    def _get_logo_path(self):
        """Return the path to the company logo image, or None if not found."""
        base_dir = os.path.dirname(os.path.abspath(__file__))
        assets_dir = os.path.abspath(os.path.join(base_dir, "..", "Data", "assets"))
        logo_path = os.path.join(assets_dir, "company_logo.png")
        if os.path.exists(logo_path):
            return logo_path
        # Try jpg fallback
        logo_path_jpg = os.path.join(assets_dir, "company_logo.jpg")
        if os.path.exists(logo_path_jpg):
            return logo_path_jpg
        return None

    def _add_logo_to_slide(self, slide, position="top-left"):
        """Add the company logo to the slide at the specified position ('top-left', 'bottom-left')."""
        logo_path = self._get_logo_path()
        if not logo_path:
            return
        if position == "bottom-left":
            left = Inches(0.2)
            top = Inches(6.5)
        else:  # top-left
            left = Inches(0.2)
            top = Inches(0.2)
        width = Inches(1.0)
        try:
            slide.shapes.add_picture(logo_path, left, top, width=width)
        except Exception as e:
            logger.warning(f"Could not add logo: {e}")

    def _add_contextual_image(self, slide, product_type_en: str, cover=False):
        """Add a contextual product image to the slide based on product type with dynamic sizing."""
        from PIL import Image
        base_dir = os.path.dirname(os.path.abspath(__file__))
        assets_dir = os.path.abspath(os.path.join(base_dir, "..", "Data", "assets"))
        
        # Convert product type to match our image naming convention
        image_name = product_type_en.lower().replace(" ", "-")
        
        # Try multiple image formats and variations
        possible_filenames = [
            f"{image_name}.jpg",
            f"{image_name}.png",
            "default_product.jpg",  # Fallback
        ]
        
        image_path = None
        for filename in possible_filenames:
            path = os.path.join(assets_dir, filename)
            if os.path.exists(path):
                logger.info(f"Found image: {filename}")
                image_path = path
                break
                
        if not image_path:
            logger.warning(f"No product image found for {product_type_en}")
            # Try general category images as fallback
            general_images = [
                "workstation.jpg",
                "computer.jpg",
                "laptop.jpg",
                "server.jpg"
            ]
            for gen_img in general_images:
                path = os.path.join(assets_dir, gen_img)
                if os.path.exists(path):
                    image_path = path
                    logger.info(f"Using fallback image: {gen_img}")
                    break
            
            if not image_path:
                return
        
        try:
            # Get image dimensions for aspect ratio preservation
            img = Image.open(image_path)
            img_width, img_height = img.size
            aspect_ratio = img_width / img_height
            
            if cover:
                # Cover image - larger and centered
                available_width = Inches(5)  # Wider for cover
                available_height = Inches(3)  # Taller for cover
                
                # Calculate dimensions preserving aspect ratio
                if aspect_ratio > 1:  # Wider than tall
                    width = available_width
                    height = width / aspect_ratio
                else:  # Taller than wide
                    height = available_height
                    width = height * aspect_ratio
                
                # Center the image horizontally
                left = Inches(4.25) - (width / 2)  # Center point is at 4.25 inches
                top = COVER_IMAGE_COVER[1]  # Use the original top position
                
                pic = slide.shapes.add_picture(image_path, left, top, width=width)
                logger.info(f"Added cover image: {os.path.basename(image_path)} ({width/Inches(1):.2f}\"x{height/Inches(1):.2f}\")")
            else:
                # Regular slide image - smaller and right-aligned
                # Check available vertical space by looking at the content
                content_shapes = [s for s in slide.shapes if hasattr(s, 'text_frame') and s.text_frame is not None]
                
                # Find bottom of content area
                bottom_of_content = 0
                for shape in content_shapes:
                    shape_bottom = shape.top + shape.height
                    if shape_bottom > bottom_of_content:
                        bottom_of_content = shape_bottom
                
                # Available height is from bottom of content to bottom of slide with margin
                slide_height = prs.slide_height if hasattr(slide, 'prs') and hasattr(slide.prs, 'slide_height') else Inches(7.5)
                available_height = slide_height - bottom_of_content - Inches(0.5)  # 0.5" margin
                
                # Available width for the image (typically on the right side)
                available_width = Inches(2.5)
                
                # Calculate dimensions preserving aspect ratio
                if aspect_ratio > 1:  # Wider than tall
                    if available_height * aspect_ratio <= available_width:
                        height = available_height
                        width = height * aspect_ratio
                    else:
                        width = available_width
                        height = width / aspect_ratio
                else:  # Taller than wide
                    if available_width / aspect_ratio <= available_height:
                        width = available_width
                        height = width / aspect_ratio
                    else:
                        height = available_height
                        width = height * aspect_ratio
                
                # Position in the bottom right
                left = prs.slide_width - width - Inches(0.5) if hasattr(slide, 'prs') else Inches(8) - width
                top = slide_height - height - Inches(0.5) if hasattr(slide, 'prs') else Inches(7) - height
                
                pic = slide.shapes.add_picture(image_path, left, top, width=width)
                logger.info(f"Added slide image: {os.path.basename(image_path)} ({width/Inches(1):.2f}\"x{height/Inches(1):.2f}\")")
        except Exception as e:
            logger.error(f"Failed to add image {image_path}: {e}", exc_info=True)

    async def generate_ppt(self, data: dict, output_path: str = "Sales_Pitch_Deck.pptx", similar_products: List[Dict[str, Any]] = None, product_name: str = None):
        """Generate a PowerPoint presentation with proper font handling"""
        resolved_language = data.get('resolved_language', 'en')
        language_resolution = data.get('language_resolution', {})
        logger.info(f"🌐 Generating presentation in resolved language: {resolved_language}")
        
        # Set fonts based on language
        if resolved_language == "ja":
            # Use Noto Sans CJK JP since we know it's available
            self.title_font = "Noto Sans CJK JP"
            self.body_font = "Noto Sans CJK JP"
            logger.info("🈁 Using Noto Sans CJK JP for Japanese text")
        else:
            self.title_font = "Calibri"
            self.body_font = "Calibri"
        
        # Get template path
        base_dir = os.path.dirname(os.path.abspath(__file__))
        template_path = os.path.join(base_dir, "..", "Data", "assets", "template.pptx")
        template_path = os.path.abspath(template_path)
        logger.info(f"Using template.pptx at {os.path.abspath(template_path)}")
        
        if os.path.exists(template_path):
            prs = Presentation(template_path)
            logger.info("✅ Template loaded successfully")
        else:
            logger.warning("⚠️ Template not found, using default presentation")
            prs = Presentation()
        
        # --- RED THEME ---
        TITLE_COLOR = RGBColor(192, 57, 43)      # Red
        ACCENT_COLOR = RGBColor(231, 76, 60)     # Lighter red
        BODY_COLOR = RGBColor(80, 80, 80)        # Dark gray for better readability
        
        # Cover slide with proper font handling
        cover_slide = prs.slides.add_slide(prs.slide_layouts[0])
        self.hide_placeholders(cover_slide)
        
        # Add cover title with improved Japanese font support
        title_box = cover_slide.shapes.add_textbox(*COVER_TITLE_BOX)
        tf = title_box.text_frame
        tf.clear()
        tf.margin_left, tf.margin_right, tf.margin_top, tf.margin_bottom = COVER_TITLE_MARGIN
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        
        # Default title text for when we haven't set it yet
        temp_title = "Sales Pitch Deck" if resolved_language != "ja" else "営業用プレゼンテーション"
        
        # Get first slide title from deck data if available
        if data.get("slides") and len(data["slides"]) > 0:
            if isinstance(data["slides"][0], dict) and "title" in data["slides"][0]:
                temp_title = data["slides"][0]["title"]
        
        # Add the text first
        p.text = temp_title
        p.font.size = COVER_TITLE_FONT_SIZE
        p.font.color.rgb = TITLE_COLOR
        
        # Apply Japanese font with multiple fallback methods
        if resolved_language == "ja":
            logger.info(f"🈁 Setting Japanese title font: {self.title_font}")
            
            # Method 1: Set font directly on the run
            for run in p.runs:
                run.font.name = self.title_font
            
            # Method 2: Set East Asian font using XML
            try:
                from pptx.oxml.ns import qn
                p.font._element.set(qn('a:eastAsiaTheme'), 'minor')
                
                # Set the font directly on each run's XML element
                for run in p.runs:
                    run_element = run._r
                    run_props = run_element.get_or_add_rPr()
                    run_props.set(qn('a:ea'), self.title_font)
            except Exception as e:
                logger.warning(f"Could not set East Asian font properties: {e}")
            
            # Method 3: Try another approach for setting font
            try:
                rPr = p._element.get_or_add_rPr()
                rPr.set('eastAsia', self.title_font)
            except Exception as e2:
                logger.warning(f"Alternative font setting also failed: {e2}")
        else:
            # For non-Japanese, simply set the font name
            p.font.name = self.title_font
        
        # Top accent line (red)
        top_line = cover_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, *COVER_TOP_LINE
        )
        top_line.fill.solid()
        top_line.fill.fore_color.rgb = ACCENT_COLOR
        top_line.line.fill.background()
        # Bottom accent line (red)
        bottom_line = cover_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, *COVER_BOTTOM_LINE
        )
        bottom_line.fill.solid()
        bottom_line.fill.fore_color.rgb = ACCENT_COLOR
        bottom_line.line.fill.background()
        # Subtitle: Quotation for {product category} from Otsuka Corporation (dynamic for Japanese)
        product_type = None
        product_type_en = None
        product_type_localized = None
        if product_name:
            cat = self.get_product_category(product_name, resolved_language)
            product_type_en = cat["en"]
            product_type_localized = cat["localized"]
        if resolved_language == "ja":
            subtitle_text = f"大塚株式会社からの{product_type_localized if product_type_localized else '製品'}見積書"
        else:
            subtitle_text = f"Quotation for {product_type_localized if product_type_localized else 'Product'} from Otsuka Co."
        subtitle_box = cover_slide.shapes.add_textbox(*COVER_SUBTITLE_BOX)
        tf_sub = subtitle_box.text_frame
        tf_sub.clear()
        tf_sub.margin_left, tf_sub.margin_right, tf_sub.margin_top, tf_sub.margin_bottom = COVER_SUBTITLE_MARGIN
        subtitle_p = tf_sub.paragraphs[0]
        self._apply_font_to_paragraph(subtitle_p, subtitle_text, is_title=False)
        subtitle_p.font.size = COVER_SUBTITLE_FONT_SIZE
        subtitle_p.font.color.rgb = ACCENT_COLOR
        subtitle_p.alignment = PP_ALIGN.CENTER
        # Remove company info box (no third text)
        # Add logo to cover slide (bottom left)
        self._add_logo_to_slide(cover_slide, position="bottom-left")
        # Add contextual image to cover slide if product_type is available (large and centered)
        if product_type_en:
            self._add_contextual_image(cover_slide, product_type_en, cover=True)
        # Content slides
        for i, slide_data in enumerate(data.get("slides", [])):
            if not isinstance(slide_data, dict):
                logger.warning(f"Slide data at index {i} is not a dict: {slide_data}")
                continue
            title = slide_data.get("title", f"Slide {i+1}")
            content = slide_data.get("content", [])
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            self.hide_placeholders(slide)
            # Title with improved Japanese font handling
            title_box = slide.shapes.add_textbox(*CONTENT_TITLE_BOX)
            tf = title_box.text_frame
            tf.clear()
            tf.margin_left, tf.margin_right, tf.margin_top, tf.margin_bottom = CONTENT_TITLE_MARGIN
            title_p = tf.paragraphs[0]
            
            # Add the text first
            title_p.text = title
            title_p.font.size = CONTENT_TITLE_FONT_SIZE
            title_p.font.bold = True
            title_p.font.color.rgb = TITLE_COLOR
            title_p.alignment = PP_ALIGN.CENTER
            
            # Apply appropriate font with better Japanese support
            if resolved_language == "ja":
                logger.info(f"🈁 Setting Japanese content title font: {self.title_font}")
                
                # Method 1: Set font directly on the runs
                for run in title_p.runs:
                    run.font.name = self.title_font
                
                # Method 2: Set East Asian font using XML
                try:
                    from pptx.oxml.ns import qn
                    title_p.font._element.set(qn('a:eastAsiaTheme'), 'minor')
                    
                    # Set the font directly on each run's XML element
                    for run in title_p.runs:
                        run_element = run._r
                        run_props = run_element.get_or_add_rPr()
                        run_props.set(qn('a:ea'), self.title_font)
                except Exception as e:
                    logger.warning(f"Could not set East Asian font properties for content title: {e}")
            else:
                # For non-Japanese, simply set the font name
                title_p.font.name = self.title_font
            # Decorative underline (red)
            underline_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, *CONTENT_UNDERLINE
            )
            underline_shape.fill.solid()
            underline_shape.fill.fore_color.rgb = ACCENT_COLOR
            underline_shape.line.fill.background()
            # Content
            content_box = slide.shapes.add_textbox(*CONTENT_BOX)
            tf = content_box.text_frame
            tf.word_wrap = True
            tf.clear()
            tf.margin_left, tf.margin_right, tf.margin_top, tf.margin_bottom = CONTENT_MARGIN
            for j, line_text in enumerate(content):
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                bullet_text = f"• {line_text}"
                self._apply_font_to_paragraph(p, bullet_text, is_title=False)
                p.font.size = CONTENT_FONT_SIZE
                p.font.color.rgb = BODY_COLOR if j % 2 == 0 else RGBColor(120, 20, 20)
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(16)
                p.space_before = Pt(8)
            # Slide number
            slide_number_box = slide.shapes.add_textbox(*CONTENT_SLIDE_NUMBER_BOX)
            tf_num = slide_number_box.text_frame
            tf_num.clear()
            num_p = tf_num.paragraphs[0]
            self._apply_font_to_paragraph(num_p, f"{i+2}", is_title=False)
            num_p.font.size = CONTENT_SLIDE_NUMBER_FONT_SIZE
            num_p.font.color.rgb = RGBColor(150, 150, 150)
            num_p.alignment = PP_ALIGN.RIGHT
            # Add logo to each content slide (top right)
            self._add_logo_to_slide(slide, position="top-left")
        # Table slides (ensure competitor/comparison slide is always present if requested)
        tables_to_process = data.get("tables", [])
        # Use localized comparison table title for detection
        comparison_title_localized = get_category_translation("comparison", resolved_language)
        def is_comparison_table(table):
            title = table.get("title", "").lower()
            return (
                "comparison" in title or "product" in title or
                comparison_title_localized.lower() in title
            )
        if similar_products or not any(is_comparison_table(t) for t in tables_to_process):
            # Always add a competitor slide if not present
            comparison_table = self.create_comparison_table_from_products(
                similar_products or [], 
                title=comparison_title_localized,
                language=resolved_language
            )
            
            if comparison_table:  # Only proceed if we got a valid table
                found = False
                for i, table_data in enumerate(tables_to_process):
                    if is_comparison_table(table_data):
                        tables_to_process[i] = comparison_table
                        found = True
                        break
                if not found:
                    tables_to_process.append(comparison_table)
        for table_data in tables_to_process:
            if not isinstance(table_data, dict):
                logger.warning(f"Table data is not a dict: {table_data}")
                continue
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            self.hide_placeholders(slide)
            # Title
            title_box = slide.shapes.add_textbox(*CONTENT_TITLE_BOX)
            tf = title_box.text_frame
            tf.clear()
            tf.margin_left, tf.margin_right, tf.margin_top, tf.margin_bottom = CONTENT_TITLE_MARGIN
            title_p = tf.paragraphs[0]
            self._apply_font_to_paragraph(title_p, table_data.get("title", "Table"), is_title=True)
            title_p.font.size = Pt(32)
            title_p.font.bold = True
            title_p.font.color.rgb = TITLE_COLOR
            title_p.alignment = PP_ALIGN.CENTER
            # Decorative underline (red)
            underline_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, *CONTENT_UNDERLINE
            )
            underline_shape.fill.solid()
            underline_shape.fill.fore_color.rgb = ACCENT_COLOR
            underline_shape.line.fill.background()
            # Table
            self.add_comparison_table(slide, table_data)
            # Add logo to each table slide (top left)
            self._add_logo_to_slide(slide, position="top-left")
        try:
            prs.save(output_path)
            logger.info("Presentation saved successfully to: %s", output_path)
            logger.info(f"🌐 Presentation generated in {resolved_language} using {language_resolution.get('method', 'unknown')} detection")
            return output_path
        except Exception as e:
            logger.error("Error saving presentation: %s", e)
            raise
