import os
import json
from openai import AzureOpenAI
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml


# Initialize Azure OpenAI client
client = AzureOpenAI(
    api_key=os.getenv("AZURE_OPENAI_API_KEY") or "6wSmubOxmIo5YWfQE45eAWnosDfSURkYmOR6yFuV3rBEULQFHW2aJQQJ99BEACYeBjFXJ3w3AAABACOG5ovQ",
    api_version="2025-01-01-preview",
    azure_endpoint="https://internship2025-teama.openai.azure.com"
)

DEPLOYMENT_NAME = "gpt-4.1-mini"

def hide_placeholders(slide):
    """Safely hide placeholders without corrupting the slide structure"""
    shapes_to_remove = []
    
    for shape in slide.shapes:
        # Check if it's a placeholder
        if hasattr(shape, 'placeholder_format') and shape.placeholder_format is not None:
            shapes_to_remove.append(shape)
        # Also remove text boxes that might be default placeholders
        elif hasattr(shape, 'text_frame') and shape.text_frame is not None:
            # Check if it contains placeholder text
            if shape.text_frame.text.strip() in ['Click to add title', 'Click to add subtitle', 'Click to add text']:
                shapes_to_remove.append(shape)
    
    # Remove identified placeholders
    for shape in shapes_to_remove:
        try:
            sp = shape._element
            sp.getparent().remove(sp)
        except:
            # If direct removal fails, try making it invisible
            try:
                shape.width = 0
                shape.height = 0
            except:
                pass

def clear_slide_content_safely(slide):
    """Safely clear slide content without corrupting the slide structure"""
    # Only remove content shapes, not structural elements
    shapes_to_remove = []
    
    for shape in slide.shapes:
        # Only remove shapes that are not essential slide structure
        if (hasattr(shape, 'shape_type') and 
            shape.shape_type in [1, 17, 14]):  # Text box, auto shape, picture
            shapes_to_remove.append(shape)
        elif hasattr(shape, 'text_frame') and shape.text_frame is not None:
            # Clear text content but keep the shape structure
            try:
                shape.text_frame.clear()
            except:
                shapes_to_remove.append(shape)
    
    # Remove non-essential shapes
    for shape in shapes_to_remove:
        try:
            sp = shape._element
            sp.getparent().remove(sp)
        except:
            # If removal fails, just hide it
            try:
                shape.width = 0
                shape.height = 0
            except:
                pass


def extract_ppt_structure(quotation: str) -> dict:
    """Use Azure OpenAI to generate a detailed and persuasive sales pitch deck structure from the quotation."""
    prompt = f"""
You are a business assistant. Based on the product quotation below, generate a structured and persuasive PowerPoint sales pitch deck in **valid JSON** format.

### QUOTATION
\"\"\"
{quotation}
\"\"\"

### TASKS
1. Analyze the quotation to identify:
   - Customer name
   - Product name
   - Specifications (CPU, RAM, Storage, etc.)
   - Price, Delivery Timeline, Warranty, Support options

2. Generate a slide deck in this order:
   1. Customer Need
   2. Our Solution
   3. Product Overview (specs)
   4. Pricing Breakdown
   5. Warranty & Support
   6. Product Comparison (see below)
   7. Delivery Timeline
   8. Call to Action

Each slide must contain a **title** and 5–6 persuasive bullet points.

3. Add a **comparison table** duplicating the same product 3 times with slightly varied names:
   - Copy the product specs from the quotation
   - Change the name to `"Product Name Variant 1"`, `"Variant 2"`, etc.
   - Add them to a table with the following structure:

### JSON OUTPUT FORMAT
Return your response as valid JSON:
{{
  "slides": [
    {{
      "title": "Slide Title",
      "content": ["Bullet 1", "Bullet 2", "..."]
    }},
    ...
  ],
  "tables": [
    {{
      "title": "Product Comparison",
      "columns": ["Product Name", "Price", "CPU", "RAM", "Storage", "Warranty", "Support"],
      "rows": [
        ["...Variant 1...", "...", "...", "...", "...", "...", "..."],
        ["...Variant 2...", "...", "...", "...", "...", "...", "..."],
        ["...Variant 3...", "...", "...", "...", "...", "...", "..."]
      ]
    }}
  ]
}}

✅ Use ONLY the product and specifications mentioned in the quotation — do NOT make up new ones.  
✅ Use slightly varied product names for realism.  
✅ Return valid JSON ONLY — no commentary, no markdown.
"""

    response = client.chat.completions.create(
        model=DEPLOYMENT_NAME,
        messages=[{"role": "user", "content": prompt}],
        temperature=0.5
    )

    raw_output = response.choices[0].message.content.strip()

    try:
        parsed = json.loads(raw_output)
        if isinstance(parsed, list):
            return {"slides": parsed}
        elif isinstance(parsed, dict) and "slides" in parsed:
            return parsed
        else:
            raise ValueError("Parsed JSON is missing expected 'slides' structure.")
    except json.JSONDecodeError:
        print("❌ GPT response was not valid JSON. Response was:\n", raw_output)
        raise


def add_comparison_table(slide, table_data):
    rows = len(table_data["rows"]) + 1  # +1 for header
    cols = len(table_data["columns"])
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)

    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Header row
    for col, header in enumerate(table_data["columns"]):
        cell = table_shape.cell(0, col)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)

    # Data rows
    for i, row in enumerate(table_data["rows"], start=1):
        for j, value in enumerate(row):
            cell = table_shape.cell(i, j)
            cell.text = str(value)  # Ensure it's a string
            cell.text_frame.paragraphs[0].font.size = Pt(12)


def generate_ppt(data: dict, output_path: str = "Sales_Pitch_Deck.pptx"):
    # Load template if it exists, otherwise create new presentation
    if os.path.exists('template.pptx'):
        print("📋 Using template.pptx")
        prs = Presentation('template.pptx')
    else:
        print("📋 Creating new presentation (no template found)")
        prs = Presentation()
    
    TITLE_FONT = "Segoe UI Semibold"
    BODY_FONT = "Segoe UI"
    TITLE_COLOR = RGBColor(44, 62, 80)
    ACCENT_COLOR = RGBColor(52, 152, 219)

    # Create cover slide
    cover_slide = prs.slides.add_slide(prs.slide_layouts[0])
    hide_placeholders(cover_slide)  # Hide template placeholders
    
    # Add cover slide content
    title_box = cover_slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(7), Inches(2))
    tf = title_box.text_frame
    tf.text = "Sales Pitch Deck"
    tf.paragraphs[0].font.size = Pt(44)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = TITLE_FONT
    tf.paragraphs[0].font.color.rgb = TITLE_COLOR
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    subtitle = tf.add_paragraph()
    subtitle.text = "Generated from Quotation"
    subtitle.font.size = Pt(24)
    subtitle.font.name = BODY_FONT
    subtitle.font.color.rgb = ACCENT_COLOR
    subtitle.alignment = PP_ALIGN.CENTER

    # Create content slides
    for slide_data in data.get("slides", []):
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Use content layout
        hide_placeholders(slide)  # Hide template placeholders
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(8), Inches(1))
        tf = title_box.text_frame
        tf.text = slide_data["title"]
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.name = TITLE_FONT
        tf.paragraphs[0].font.color.rgb = TITLE_COLOR

        # Add accent line
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(0.1), Inches(4))
        line.fill.solid()
        line.fill.fore_color.rgb = ACCENT_COLOR
        line.line.fill.background()

        # Add content
        content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8), Inches(5))
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, line_text in enumerate(slide_data["content"]):
            if i == 0:
                # First paragraph already exists
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = f"• {line_text}"
            p.font.size = Pt(20)
            p.font.name = BODY_FONT
            p.font.color.rgb = RGBColor(80, 80, 80)
            p.space_after = Pt(12)

    # Create table slides
    for table_data in data.get("tables", []):
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Use content layout
        hide_placeholders(slide)  # Hide template placeholders
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(8), Inches(1))
        tf = title_box.text_frame
        tf.text = table_data["title"]
        tf.paragraphs[0].font.size = Pt(32)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.name = TITLE_FONT
        tf.paragraphs[0].font.color.rgb = TITLE_COLOR
        
        # Add table
        add_comparison_table(slide, table_data)

    # Save presentation
    try:
        prs.save(output_path)
        print(f"✅ Presentation saved successfully to: {output_path}")
    except Exception as e:
        print(f"❌ Error saving presentation: {e}")
        raise


if __name__ == "__main__":
    quotation_input = """
    Here is the formal quote for one unit of the Workstation Pro Professional:

    Customer: [Your Company Name]
    Product: Workstation Pro Professional

    CPU: Intel Xeon W-2295
    RAM: 32GB DDR4
    Storage: 1TB NVMe SSD
    Quantity: 1 unit
    Unit Price: $3,499.99
    Delivery Timeline: 2-3 weeks
    Warranty: 1-year included
    Support: Optional setup and integration
    """

    try:
        print("🧠 Generating structured pitch deck from quotation...")
        structured_data = extract_ppt_structure(quotation_input)

        print("🖼️ Creating PowerPoint presentation...")
        generate_ppt(structured_data)

    except Exception as e:
        print("❌ Error:", e)